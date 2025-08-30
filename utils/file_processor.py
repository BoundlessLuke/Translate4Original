import os
import docx
import openpyxl
import pptx
from pptx.util import Inches
import tempfile

# 处理DOCX文件
def process_docx(file_path, source_lang, target_lang, translate_func, custom_prompt=None):
    # 打开文档
    doc = docx.Document(file_path)
    
    # 遍历所有段落
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            # 翻译段落文本
            translated_text = translate_func(paragraph.text, source_lang, target_lang, custom_prompt)
            # 保留原格式，只替换文本
            original_runs = paragraph.runs
            if original_runs:
                # 清空段落
                paragraph.clear()
                # 添加翻译后的文本，保留第一个run的所有格式
                translated_run = paragraph.add_run(translated_text)
                translated_run.font.name = original_runs[0].font.name
                translated_run.font.size = original_runs[0].font.size
                translated_run.font.bold = original_runs[0].font.bold
                translated_run.font.italic = original_runs[0].font.italic
                translated_run.font.underline = original_runs[0].font.underline
                translated_run.font.color.rgb = original_runs[0].font.color.rgb if hasattr(original_runs[0].font.color, 'rgb') else None
                # 保留段落的行距和对齐方式
                paragraph.paragraph_format.line_spacing = paragraph.paragraph_format.line_spacing
                paragraph.paragraph_format.alignment = paragraph.paragraph_format.alignment
    
    # 遍历所有表格
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    # 翻译表格单元格文本
                    translated_text = translate_func(cell.text, source_lang, target_lang, custom_prompt)
                    # 保存单元格的格式
                    original_font = None
                    original_size = None
                    if cell.paragraphs and cell.paragraphs[0].runs:
                        original_font = cell.paragraphs[0].runs[0].font.name
                        original_size = cell.paragraphs[0].runs[0].font.size
                    # 设置翻译后的文本
                    cell.text = translated_text
                    # 恢复单元格的格式
                    if cell.paragraphs and cell.paragraphs[0].runs:
                        cell.paragraphs[0].runs[0].font.name = original_font
                        cell.paragraphs[0].runs[0].font.size = original_size
    
    return doc, 'docx'

# 处理XLSX文件
def process_xlsx(file_path, source_lang, target_lang, translate_func, custom_prompt=None):
    # 打开工作簿
    wb = openpyxl.load_workbook(file_path)
    
    # 遍历所有工作表
    for sheet in wb.worksheets:
        # 遍历所有单元格
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.strip():
                    # 翻译单元格文本
                    try:
                        translated_text = translate_func(cell.value, source_lang, target_lang, custom_prompt)
                        cell.value = translated_text
                    except:
                        # 翻译失败时保留原文本
                        pass
    
    return wb, 'xlsx'

# 处理PPTX文件
def process_pptx(file_path, source_lang, target_lang, translate_func, custom_prompt=None):
    # 打开演示文稿
    prs = pptx.Presentation(file_path)
    
    # 遍历所有幻灯片
    for slide in prs.slides:
        # 遍历所有形状
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # 保存原始字体属性
                original_text_frame = shape.text_frame
                original_font = None
                original_size = None
                original_color = None
                
                if original_text_frame and original_text_frame.paragraphs:
                    for para in original_text_frame.paragraphs:
                        if para.runs:
                            original_font = para.runs[0].font.name
                            original_size = para.runs[0].font.size
                            original_color = para.runs[0].font.color.rgb if hasattr(para.runs[0].font.color, 'rgb') else None
                            break
                
                # 翻译形状文本
                translated_text = translate_func(shape.text, source_lang, target_lang, custom_prompt)
                shape.text = translated_text
                
                # 恢复字体属性
                if original_text_frame and original_text_frame.paragraphs and original_font:
                    for para in original_text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = original_font
                            if original_size:
                                run.font.size = original_size
                            if original_color:
                                run.font.color.rgb = original_color
            
            # 处理表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            # 保存原始字体属性
                            original_font = None
                            original_size = None
                            if cell.text_frame and cell.text_frame.paragraphs:
                                for para in cell.text_frame.paragraphs:
                                    if para.runs:
                                        original_font = para.runs[0].font.name
                                        original_size = para.runs[0].font.size
                                        break
                            
                            translated_text = translate_func(cell.text, source_lang, target_lang, custom_prompt)
                            cell.text = translated_text
                            
                            # 恢复字体属性
                            if cell.text_frame and cell.text_frame.paragraphs and original_font:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.name = original_font
                                        if original_size:
                                            run.font.size = original_size
    
    return prs, 'pptx'

# 处理DOC文件（转换为DOCX后处理）
def process_doc(file_path, source_lang, target_lang, translate_func, custom_prompt=None):
    # 这里简化处理，实际上可能需要使用python-docx2txt或其他库
    # 或者提示用户将DOC文件转换为DOCX后再上传
    # 为了演示，我们创建一个新的DOCX文件
    doc = docx.Document()
    doc.add_heading('DOC文件翻译提示', 0)
    doc.add_paragraph('由于技术限制，.doc文件需要先转换为.docx格式后再进行翻译。')
    doc.add_paragraph('请使用Microsoft Word或其他工具将文件转换后重新上传。')
    return doc, 'docx'

# 处理XLS文件（转换为XLSX后处理）
def process_xls(file_path, source_lang, target_lang, translate_func, custom_prompt=None):
    # 类似DOC文件的处理方式
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = '翻译提示'
    ws['A1'] = 'XLS文件翻译提示'
    ws['A2'] = '由于技术限制，.xls文件需要先转换为.xlsx格式后再进行翻译。'
    ws['A3'] = '请使用Microsoft Excel或其他工具将文件转换后重新上传。'
    return wb, 'xlsx'

# 根据文件类型选择相应的处理函数
def process_file(file_path, source_lang, target_lang, translate_func, custom_prompt=None):
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if file_ext == '.docx':
        return process_docx(file_path, source_lang, target_lang, translate_func, custom_prompt)
    elif file_ext == '.doc':
        return process_doc(file_path, source_lang, target_lang, translate_func, custom_prompt)
    elif file_ext == '.xlsx':
        return process_xlsx(file_path, source_lang, target_lang, translate_func, custom_prompt)
    elif file_ext == '.xls':
        return process_xls(file_path, source_lang, target_lang, translate_func, custom_prompt)
    elif file_ext == '.pptx':
        return process_pptx(file_path, source_lang, target_lang, translate_func, custom_prompt)
    elif file_ext == '.ppt':
        # 对于PPT文件，创建一个新的PPTX文件作为提示
        prs = pptx.Presentation()
        slide_layout = prs.slide_layouts[0]  # 使用标题幻灯片布局
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "PPT文件翻译提示"
        subtitle.text = "由于技术限制，.ppt文件需要先转换为.pptx格式后再进行翻译。\n请使用Microsoft PowerPoint或其他工具将文件转换后重新上传。"
        return prs, 'pptx'
    else:
        raise ValueError(f"不支持的文件格式: {file_ext}")

# 保存翻译后的文件
def save_translated_file(content, file_type, output_path):
    if file_type == 'docx':
        content.save(output_path)
    elif file_type == 'xlsx':
        content.save(output_path)
    elif file_type == 'pptx':
        content.save(output_path)
    else:
        raise ValueError(f"不支持的文件类型: {file_type}")
